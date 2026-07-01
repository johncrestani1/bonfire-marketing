#!/usr/bin/env python3
r"""
Bonfire Terminal 18-Page Sales Brochure Builder — V6
Format: 7" x 9" portrait, high-res print
V6 changes from v5:
  - P3 TOC: Updated to match actual page titles
  - P4: Section title matches TOC
  - P10: Added page title
  - P15: Redesigned pricing — $5k full-width top, $10k/$18k side-by-side below
  - P15: Rectangular JC McLaren photo (not circle headshot)
  - P16: Fixed AS SEEN ON logo overlap — 5 logos, constrained width
  - P16: Section title matches TOC
  - P17: Redesigned JC section — rectangular McLaren photo left, text right
  - P17: FAQ #7 updated — ChatGPT -> OpenClaw differentiation
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path
from PIL import Image as PILImage

SLIDE_W = Inches(7); SLIDE_H = Inches(9)
RED = RGBColor(0xF9,0x42,0x38); GOLD = RGBColor(0xB2,0x9D,0x60)
LGOLD = RGBColor(0xD0,0xB8,0x71); DARK = RGBColor(0x1A,0x1A,0x1A)
WHITE = RGBColor(0xFF,0xFF,0xFF); OFFWH = RGBColor(0xF5,0xF5,0xF5)
CHAR = RGBColor(0x2A,0x2A,0x2A); DKRED = RGBColor(0xC0,0x30,0x2B)
GREEN = RGBColor(0x00,0xCC,0x66); GRAY = RGBColor(0x80,0x80,0x80)
DKGOLD = RGBColor(0x8A,0x7A,0x4A)
GOLDBG = RGBColor(0xF5,0xF0,0xE0); GOLDTINT = RGBColor(0xEE,0xE6,0xD0)
NAVY = RGBColor(0x01,0x31,0x61)
BASE = Path("D:/marketing/18pg-sales-brochure"); IMG = BASE/"images"
OUT = BASE/"Bonfire_Terminal_Brochure_7x9_v7.pptx"
CIRC = IMG/"circular-headshots"
LOGO_PATH = str(IMG/"bonfire-logo-128.png")
LOGO_LG = str(IMG/"bonfire-logo-512.png")
JC_MCLAREN = str(IMG/"john-mclaren-p15.jpg")
JC_STUDENT = str(IMG/"john-student-p17.jpg")
COMMERCE_TAB = str(IMG/"commerce-tab-cropped.png")
JC_BOOKCOVER = str(IMG/"book-cover-p2.jpg")
JC_SIGNATURE = str(IMG/"jc-signature-gold.png")
JC_SIG_WHITE = str(IMG/"jc-signature-white.png")
CHECKGREEN = RGBColor(0x25,0x9F,0x70)

# Print-safe boundaries
HEADER_BOTTOM = Inches(0.59)  # Header ends here
FOOTER_TOP = Inches(8.7)       # Footer starts here
MAX_CONTENT_Y = Inches(8.4)    # Last safe Y for content
HERO_H = Inches(2.0)           # Fixed hero image height
HERO_BOTTOM = Inches(2.59)     # Header bottom + hero height
CONTENT_START = Inches(2.7)    # Where content begins after hero

MX = Inches(0.4); CX = MX; CW = Inches(6.2); CR = Inches(6.6)

def img(sub, pat):
    d = IMG/sub
    if not d.exists(): return None
    for f in sorted(d.iterdir()):
        if pat.lower() in f.name.lower() and f.suffix.lower() in ('.png','.jpg','.jpeg'):
            return str(f)
    return None

def logo(name):
    # bonfire-sot first (actual logos from Bonfire source code), then fallback to other sets
    for sub in ['bonfire-sot','logo-dev','tech-logos-png','devicon-logos','logo-badges','wikimedia-logos','tech-logos']:
        d = IMG/sub
        if not d.exists(): continue
        for f in sorted(d.iterdir()):
            if name.lower() in f.name.lower() and f.suffix.lower() in ('.png','.jpg','.jpeg'):
                return str(f)
    return None

def rect(s,l,t,w,h,fill,line=None):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,l,t,w,h)
    sh.fill.solid(); sh.fill.fore_color.rgb=fill
    if line: sh.line.color.rgb=line; sh.line.width=Pt(1)
    else: sh.line.fill.background()
    return sh

def txt(s,l,t,w,h,text,sz=12,clr=WHITE,bold=False,align=PP_ALIGN.LEFT,font='Inter',italic=False):
    tb=s.shapes.add_textbox(l,t,w,h); tf=tb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.text=text; p.font.size=Pt(sz); p.font.color.rgb=clr
    p.font.bold=bold; p.font.name=font; p.alignment=align; p.font.italic=italic
    return tb

def pic(s,path,l,t,w=None,h=None):
    if not path or not os.path.exists(path) or path.lower().endswith('.svg'): return None
    try:
        kw={}
        if w: kw['width']=w
        if h: kw['height']=h
        return s.shapes.add_picture(path,l,t,**kw)
    except Exception as e:
        print(f"  WARN pic: {path}: {e}"); return None

def oval(s,l,t,w,h,fill,line=None):
    sh=s.shapes.add_shape(MSO_SHAPE.OVAL,l,t,w,h)
    sh.fill.solid(); sh.fill.fore_color.rgb=fill
    if line: sh.line.color.rgb=line; sh.line.width=Pt(1.5)
    else: sh.line.fill.background()
    return sh

def bg(s,c): s.background.fill.solid(); s.background.fill.fore_color.rgb=c

def header(s):
    """V5 header: Logo + AI | BONFIRE TERMINAL"""
    rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.36), DARK)
    rect(s, Inches(0), Inches(0.36), SLIDE_W, Inches(0.19), CHAR)
    rect(s, Inches(0), Inches(0.55), SLIDE_W, Inches(0.04), GOLD)
    if os.path.exists(LOGO_PATH):
        pic(s, LOGO_PATH, Inches(0.15), Inches(0.04), w=Inches(0.28), h=Inches(0.28))
    txt(s, Inches(0.5), Inches(0.05), Inches(0.4), Inches(0.28),
        "AI", sz=12, clr=GOLD, bold=True, font='Orbitron')
    txt(s, Inches(0.88), Inches(0.05), Inches(0.15), Inches(0.28),
        "|", sz=12, clr=LGOLD, bold=False, font='Orbitron')
    txt(s, Inches(1.05), Inches(0.05), Inches(3.5), Inches(0.28),
        "BONFIRE TERMINAL", sz=11, clr=RED, bold=True, font='Orbitron')

def footer(s, page_num=None):
    rect(s, Inches(0), Inches(8.7), SLIDE_W, Inches(0.3), DARK)
    if page_num:
        txt(s, Inches(5.8), Inches(8.72), Inches(0.8), Inches(0.25),
            str(page_num), sz=8, clr=GOLD, font='Orbitron', align=PP_ALIGN.RIGHT)

def section_num(s, num, y=None):
    """Gold corner badge (Regal-inspired) at bottom-left of hero area."""
    if y is None:
        y = HERO_BOTTOM - Inches(0.1)
    tri = s.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                             Inches(0), y - Inches(0.6), Inches(1.0), Inches(0.7))
    tri.fill.solid(); tri.fill.fore_color.rgb = GOLD; tri.line.fill.background()
    tri.rotation = 180.0
    rect(s, Inches(0), y - Inches(0.05), Inches(0.7), Inches(0.55), DARK)
    rect(s, Inches(0), y - Inches(0.05), Inches(0.04), Inches(0.55), GOLD)
    txt(s, Inches(0.06), y - Inches(0.02), Inches(0.6), Inches(0.5),
        str(num), sz=36, clr=GOLD, bold=True, font='Orbitron', align=PP_ALIGN.CENTER)

def hero_img(s, pattern, top=HEADER_BOTTOM, max_h=None):
    """Place hero image with HARD CUTOFF at max_h. Uses pre-cropped *-hero.jpg if available."""
    # Prefer pre-cropped hero version (correct aspect ratio, no stretching)
    hero_path = img("ai-characters", pattern + "-hero")
    path = hero_path or img("ai-characters", pattern)
    if not path: return top
    if max_h is None:
        max_h = HERO_H
    try:
        im = PILImage.open(path)
        w, h = im.size
        im.close()
        aspect = h / w
        target_w = SLIDE_W
        calc_h = int(target_w * aspect)
        if calc_h <= max_h:
            # Image fits within max_h — use natural dimensions (no stretching)
            pic(s, path, Inches(0), top, w=target_w, h=calc_h)
            return top + calc_h
        else:
            # Image too tall — place at max_h (will crop bottom visually)
            pic(s, path, Inches(0), top, w=target_w, h=max_h)
            return top + max_h
    except:
        pic(s, path, Inches(0), top, w=SLIDE_W, h=max_h)
        return top + max_h

def logo_row(s, names, left, top, icon_size=0.42, gap=0.58, max_count=10):
    placed = 0
    for name in names:
        if placed >= max_count: break
        path = logo(name)
        if not path: continue
        x = left + Inches(placed * gap)
        pic(s, path, x, top, w=Inches(icon_size), h=Inches(icon_size))
        placed += 1
    return placed

def logo_grid_row(s, category_name, desc, names, left, top, row_width, row_height,
                  icon_size=0.4):
    """Logo grid row with category description paragraph."""
    rect(s, left, top, row_width, Inches(0.28), DARK)
    txt(s, left+Inches(0.1), top+Inches(0.02), row_width-Inches(0.2), Inches(0.24),
        category_name, sz=9, clr=GOLD, bold=True, font='Orbitron')
    desc_top = top + Inches(0.28)
    if desc:
        rect(s, left, desc_top, row_width, Inches(0.22), GOLDBG)
        txt(s, left+Inches(0.1), desc_top+Inches(0.02), row_width-Inches(0.2), Inches(0.18),
            desc, sz=7, clr=DKGOLD, font='Inter', italic=True)
        content_top = desc_top + Inches(0.22)
    else:
        content_top = desc_top
    content_h = row_height - (content_top - top)
    rect(s, left, content_top, row_width, content_h, WHITE, line=GOLD)
    placed = 0
    max_per_row = min(len(names), 8)
    gap = (row_width - Inches(0.2)) / max(max_per_row, 1)
    for i, name in enumerate(names):
        if placed >= max_per_row: break
        path = logo(name)
        if not path: continue
        x = left + Inches(0.15) + gap * placed
        y = content_top + Inches(0.1)
        pic(s, path, x, y, w=Inches(icon_size), h=Inches(icon_size))
        # Slash-format labels for AI providers (product/brand)
        ai_labels = {"claude-color":"Claude/Anthropic","openai":"OpenAI",
                     "gemini-color":"Gemini/Google","xai-grok":"Grok/xAI",
                     "meta-llama":"Llama/Meta"}
        clean_name = ai_labels.get(name,
            name.replace("-color","").replace("-logo","").replace("-"," ").title())
        if len(clean_name) > 16: clean_name = clean_name[:16]
        txt(s, x - Inches(0.05), y + Inches(icon_size + 0.02),
            Inches(icon_size + 0.1), Inches(0.2),
            clean_name, sz=5, clr=CHAR, font='Inter', align=PP_ALIGN.CENTER)
        placed += 1
    return placed

def generate_aide_badge():
    badge_path = str(IMG / "certification" / "aide-certified-badge.png")
    if os.path.exists(badge_path):
        return badge_path
    from PIL import ImageDraw, ImageFont
    size = 300
    im = PILImage.new('RGBA', (size, size), (0, 0, 0, 0))
    draw = ImageDraw.Draw(im)
    draw.ellipse([5, 5, size-5, size-5], fill=(26, 26, 26, 255))
    draw.ellipse([20, 20, size-20, size-20], fill=(255, 255, 255, 255))
    draw.ellipse([15, 15, size-15, size-15], outline=(249, 66, 56, 255), width=4)
    cx, cy = size//2, size//2 - 15
    flame_pts = [(cx, cy-45), (cx+25, cy+10), (cx+15, cy+35),
                 (cx, cy+20), (cx-15, cy+35), (cx-25, cy+10)]
    draw.polygon(flame_pts, fill=(249, 66, 56, 255))
    inner_pts = [(cx, cy-20), (cx+10, cy+5), (cx+5, cy+20),
                 (cx, cy+12), (cx-5, cy+20), (cx-10, cy+5)]
    draw.polygon(inner_pts, fill=(255, 180, 50, 255))
    try:
        font_sm = ImageFont.truetype("arial.ttf", 22)
        font_lg = ImageFont.truetype("arialbd.ttf", 36)
    except:
        font_sm = ImageFont.load_default()
        font_lg = font_sm
    draw.text((cx, cy+55), "CERTIFIED", fill=(26, 26, 26, 255), font=font_sm, anchor="mt")
    draw.text((cx, cy+85), "AIDE", fill=(26, 26, 26, 255), font=font_lg, anchor="mt")
    im.save(badge_path)
    print(f"  Generated AIDE badge: {badge_path}")
    return badge_path


def build():
    prs = Presentation(); prs.slide_width = SLIDE_W; prs.slide_height = SLIDE_H
    BL = prs.slide_layouts[6]

    # ============================================================
    # P1: COVER
    # ============================================================
    print("P1: Cover")
    s = prs.slides.add_slide(BL); bg(s, DARK)
    rect(s, Inches(0.25), Inches(0.25), Inches(6.5), Inches(8.5), DARK, line=GOLD)
    rect(s, Inches(0.35), Inches(0.35), Inches(6.3), Inches(8.3),
         DARK, line=RGBColor(0x3A,0x35,0x28))
    if os.path.exists(LOGO_LG):
        pic(s, LOGO_LG, Inches(2.25), Inches(1.8), w=Inches(2.5), h=Inches(2.5))
    txt(s, Inches(0.5), Inches(4.6), Inches(6), Inches(0.8), "BONFIRE",
        sz=42, clr=WHITE, bold=True, font='Orbitron', align=PP_ALIGN.CENTER)
    txt(s, Inches(0.5), Inches(5.3), Inches(6), Inches(0.6), "TERMINAL",
        sz=42, clr=WHITE, bold=True, font='Orbitron', align=PP_ALIGN.CENTER)

    # ============================================================
    # P2: QUOTES
    # ============================================================
    print("P2: Authority Quotes")
    s = prs.slides.add_slide(BL); bg(s, WHITE)
    header(s); footer(s, 1)
    quotes = [
        ("isaac-newton", "Isaac Newton",
         "\u201CIf I have seen further, it is by standing on the shoulders of giants.\u201D",
         "Every AI breakthrough stands on centuries of scientific progress. The lineage from Newton\u2019s calculus to modern neural networks is unbroken."),
        ("jeff-bezos", "Jeff Bezos",
         "\u201CWe are at the beginning of a golden age of AI. Recent advancements have already led to invention that previously lived in the realm of science fiction.\u201D",
         "Founder of Amazon. Early investor in artificial intelligence at scale."),
        ("elon-musk", "Elon Musk",
         "\u201CThe pace of progress in artificial intelligence is incredibly fast. It is growing at a pace close to exponential. The risk of something seriously dangerous happening is in the five-year timeframe.\u201D",
         "CEO of Tesla, SpaceX, and xAI. Building Grok and investing billions in AI compute."),
    ]
    y = 0.75
    for circ_name, name, quote, note in quotes:
        circ_img = str(CIRC / f"{circ_name}-circle.png")
        if os.path.exists(circ_img):
            oval(s, CX, Inches(y + 0.05), Inches(0.7), Inches(0.7), WHITE, line=GOLD)
            pic(s, circ_img, CX + Inches(0.04), Inches(y + 0.09), w=Inches(0.62), h=Inches(0.62))
        txt(s, CX + Inches(0.85), Inches(y - 0.02), Inches(4.9), Inches(0.25),
            name.upper(), sz=10, clr=RED, bold=True, font='Orbitron')
        txt(s, CX + Inches(0.85), Inches(y + 0.25), Inches(4.9), Inches(0.7), quote,
            sz=8, clr=DARK, font='Inter', italic=True)
        txt(s, CX + Inches(0.85), Inches(y + 0.65), Inches(4.9), Inches(0.2),
            note, sz=6, clr=GRAY, font='Inter')
        rect(s, CX + Inches(0.85), Inches(y + 0.9), Inches(4.9), Inches(0.005), LGOLD)
        y += 1.1

    # Meet John Crestani — v6.8: taller image, all text beside photo, white sig
    jc_note_y = Inches(4.05)
    rect(s, CX, jc_note_y, CW, Inches(3.65), DARK)
    if os.path.exists(JC_BOOKCOVER):
        pic(s, JC_BOOKCOVER, CX, jc_note_y, w=Inches(2.8), h=Inches(2.35))
    txt(s, CX+Inches(2.9), jc_note_y+Inches(0.08), Inches(3.2), Inches(0.25),
        "MEET JOHN CRESTANI", sz=10, clr=GOLD, bold=True, font='Orbitron')
    txt(s, CX+Inches(2.9), jc_note_y+Inches(0.38), Inches(3.2), Inches(3.0),
        "Author of Work At Home Secrets (20,000+ copies sold via TV infomercials). "
        "John went from asking his boss for a raise to earning $500K/month selling "
        "online. His students have gone on to build six and seven-figure businesses.\n\n"
        "\u201CThose that control AI will control the world. Those that don\u2019t control "
        "AI will be the ones controlled. We work together for freedom, family, "
        "and fortune.\u201D\n\n"
        "In the AI age, there is immense opportunity to grow and prosper. "
        "Bonfire Terminal puts the same decisive advantage that Fortune 500 companies "
        "have on YOUR desktop. The people who have the best success are naturally "
        "resourceful, willing to invest time into learning, and are very proactive.",
        sz=7, clr=WHITE, font='Inter')
    if os.path.exists(JC_SIG_WHITE):
        pic(s, JC_SIG_WHITE, CX+Inches(4.5), jc_note_y+Inches(3.15), w=Inches(1.5))

    # ============================================================
    # P3: TOC — V6: Updated to match actual page titles
    # ============================================================
    print("P3: Table of Contents")
    s = prs.slides.add_slide(BL); bg(s, WHITE)
    header(s); footer(s, 2)
    txt(s, CX, Inches(2.5), CW, Inches(0.6), "TABLE OF CONTENTS",
        sz=24, clr=DARK, bold=True, font='Orbitron', align=PP_ALIGN.CENTER)
    rect(s, Inches(2.5), Inches(3.2), Inches(2), Inches(0.03), GOLD)
    secs = [
        ("1.", "An Introduction to AI", "Pages 3\u20136"),
        ("2.", "The Bonfire Platform", "Pages 7\u201311"),
        ("3.", "Why Choose Bonfire Terminal", "Pages 12\u201317"),
    ]
    y = 4.0
    for num, title, pages in secs:
        txt(s, Inches(1.5), Inches(y), Inches(0.5), Inches(0.3),
            num, sz=14, clr=GOLD, bold=True, font='Orbitron', align=PP_ALIGN.RIGHT)
        txt(s, Inches(2.2), Inches(y), Inches(3.0), Inches(0.3),
            title, sz=14, clr=DARK, bold=True, font='Inter')
        txt(s, Inches(5.0), Inches(y), Inches(1.2), Inches(0.3),
            pages, sz=10, clr=GOLD, font='Inter', align=PP_ALIGN.RIGHT)
        if y < 5.2:
            rect(s, Inches(1.5), Inches(y + 0.4), Inches(4.7), Inches(0.005), LGOLD)
        y += 0.7

    # ============================================================
    # P4: SECTION 1 OPENER
    # ============================================================
    print("P4: Introduction to AI")
    s = prs.slides.add_slide(BL); bg(s, WHITE)
    header(s); footer(s, 3)
    hero_bottom = hero_img(s, "eniac-classic", top=HEADER_BOTTOM, max_h=HERO_H)
    txt(s, CX, hero_bottom + Inches(0.05), CW, Inches(0.35),
        "AN INTRODUCTION TO ARTIFICIAL INTELLIGENCE",
        sz=13, clr=DARK, bold=True, font='Orbitron')
    col_w = Inches(2.9)
    text_top = hero_bottom + Inches(0.4)
    txt(s, CX, text_top, col_w, Inches(5.0),
        "Since the dawn of civilization, humans have sought to extend "
        "their cognitive abilities through tools. From the abacus to the "
        "printing press, each leap amplified what one person can accomplish. "
        "The first calculating machines appeared in the seventeenth century, "
        "and by the early twentieth century, mechanical computers were solving "
        "equations that would take humans years to complete by hand.\n\n"
        "The computer revolution that began with ENIAC in 1946 gave us the "
        "power to process millions of calculations per second. But it still "
        "required humans to write every instruction, debug every error, and "
        "maintain every system. The gap between human intent and machine "
        "execution remained vast.\n\n"
        "Artificial Intelligence changes this equation entirely. For the first "
        "time in history, software can understand natural language, generate "
        "original content, reason about complex problems, and make decisions "
        "\u2014 all without explicit programming. The implications for productivity, "
        "creativity, and wealth creation are staggering.",
        sz=9, clr=DARK, font='Inter')
    txt(s, CX + Inches(3.2), text_top, col_w, Inches(5.0),
        "The modern AI era began in 2017 with the Transformer architecture, "
        "published in Google\u2019s landmark paper \u2018Attention Is All You Need.\u2019 "
        "This single breakthrough enabled ChatGPT, Claude, Gemini, and every "
        "major language model in existence today. Within five years, AI went "
        "from a research curiosity to the fastest-growing technology category "
        "in the world.\n\n"
        "But the companies building these models are not building them for you. "
        "They are building platforms designed to extract maximum recurring revenue "
        "from your dependence. Every month, they raise prices, change terms, and "
        "restrict access to features you rely on.\n\n"
        "At Bletchley Park in 1943, Alan Turing and the Colossus machine proved "
        "that organized intelligence processed faster than an enemy could respond. "
        "Today, Bonfire Terminal puts that same decisive advantage "
        "on YOUR desktop \u2014 with no recurring fees, no data harvesting, and no "
        "dependency on any single provider.",
        sz=9, clr=DARK, font='Inter')

    # ============================================================
    # P5: URGENCY
    # ============================================================
    print("P5: The Cost of Waiting")
    s = prs.slides.add_slide(BL); bg(s, WHITE)
    header(s); footer(s, 4)
    txt(s, CX, Inches(0.7), CW, Inches(0.4), "THE COST OF WAITING",
        sz=18, clr=DARK, bold=True, font='Orbitron')
    txt(s, CX, Inches(1.15), CW, Inches(0.25),
        "Why AI is the greatest wealth-creation opportunity since the internet.",
        sz=9, clr=CHAR, font='Inter')
    rect(s, CX, Inches(1.6), Inches(2.8), Inches(0.4), CHAR)
    txt(s, CX, Inches(1.62), Inches(2.8), Inches(0.35), "TRADITIONAL",
        sz=11, clr=WHITE, bold=True, font='Orbitron', align=PP_ALIGN.CENTER)
    rect(s, Inches(3.5), Inches(1.6), Inches(3.0), Inches(0.4), GOLD)
    txt(s, Inches(3.5), Inches(1.62), Inches(3.0), Inches(0.35), "A.I. GHOST AGENCY",
        sz=11, clr=DARK, bold=True, font='Orbitron', align=PP_ALIGN.CENTER)
    trad = ["1 person","40 hours / week","1 client at a time","Limited by hours","$60,000 / year"]
    aipw = ["1 person","10 hours / week","50 clients simultaneously","Limited only by imagination","$300,000+ / year"]
    y = 2.1
    for i, (t, a) in enumerate(zip(trad, aipw)):
        bgc = OFFWH if i%2==0 else WHITE; gbg = GOLDBG if i%2==0 else WHITE
        rect(s, CX, Inches(y), Inches(2.8), Inches(0.45), bgc)
        rect(s, Inches(3.5), Inches(y), Inches(3.0), Inches(0.45), gbg)
        txt(s, CX+Inches(0.1), Inches(y+0.06), Inches(2.6), Inches(0.33), t,
            sz=10, clr=CHAR, font='Inter', align=PP_ALIGN.CENTER)
        c = DKGOLD if "300" in a else DARK; b = "300" in a or "50" in a
        txt(s, Inches(3.6), Inches(y+0.06), Inches(2.8), Inches(0.33), a,
            sz=10, clr=c, bold=b, font='Inter', align=PP_ALIGN.CENTER)
        y += 0.45
    rect(s, CX, Inches(4.5), CW, Inches(0.55), DARK)
    txt(s, CX+Inches(0.1), Inches(4.55), Inches(6.0), Inches(0.45),
        "\u201CAI doesn\u2019t replace you. It multiplies you.\u201D",
        sz=14, clr=GOLD, bold=True, font='Inter', align=PP_ALIGN.CENTER)
    hero_img(s, "agent-smith", top=Inches(5.3), max_h=Inches(3.0))

    # ============================================================
    # P6: CHART
    # ============================================================
    print("P6: AI Timeline Chart")
    s = prs.slides.add_slide(BL); bg(s, WHITE)
    header(s); footer(s, 5)
    txt(s, CX, Inches(0.7), CW, Inches(0.35),
        "AI MARKET GROWTH TRAJECTORY", sz=16, clr=DARK, bold=True, font='Orbitron')
    txt(s, CX, Inches(1.1), CW, Inches(0.25),
        "Global AI Market Size (Billions USD)", sz=9, clr=CHAR, font='Inter')
    chart_left = CX+Inches(0.5); chart_top = Inches(1.5)
    chart_w = Inches(5.4); chart_h = Inches(5.0)
    rect(s, chart_left, chart_top, chart_w, chart_h, WHITE, line=DARK)
    for i in range(6):
        gy = chart_top + Inches(i * 1.0)
        rect(s, chart_left, gy, chart_w, Inches(0.005), RGBColor(0xE0,0xE0,0xE0))
    y_vals = ["$300B","$250B","$200B","$150B","$100B","$50B","$0"]
    for i, v in enumerate(y_vals):
        txt(s, CX, chart_top+Inches(i*0.83)-Inches(0.1), Inches(0.55), Inches(0.2),
            v, sz=6, clr=GRAY, font='Inter', align=PP_ALIGN.RIGHT)
    points = [(0.0,4.8),(0.7,4.7),(1.4,4.5),(2.1,4.1),(2.8,3.5),(3.5,2.7),(4.2,1.7),(4.9,0.5)]
    for i in range(len(points)-1):
        x1,y1=points[i]; x2,y2=points[i+1]
        rect(s, chart_left+Inches(x1), chart_top+Inches(y1),
             Inches(x2-x1+0.05), Inches(0.06), GOLD)
    milestones = [
        (0.1,4.5,"2015","$3.2B","Deep learning boom"),
        (0.8,4.4,"2017","$12B","Transformer architecture"),
        (1.5,4.1,"2019","$27B","GPT-2 language gen"),
        (2.2,3.7,"2020","$51B","GPT-3: 175B params"),
        (2.9,3.0,"2022","$93B","ChatGPT: 100M users"),
        (3.6,2.1,"2024","$184B","GPT-4, Claude, Gemini"),
        (4.3,0.9,"2026","$280B+","BONFIRE TERMINAL"),
    ]
    for mx,my,year,val,desc in milestones:
        oval(s, chart_left+Inches(mx+0.3), chart_top+Inches(my-0.05), Inches(0.12), Inches(0.12), RED)
        clr = RED if "BONFIRE" in desc else DARK
        txt(s, chart_left+Inches(mx), chart_top+Inches(my-0.5), Inches(1.0), Inches(0.15),
            f"{year}: {val}", sz=7, clr=clr, bold=True, font='Inter')
        txt(s, chart_left+Inches(mx), chart_top+Inches(my-0.33), Inches(1.0), Inches(0.3),
            desc, sz=5, clr=CHAR, font='Inter')
    txt(s, Inches(4.5), Inches(6.7), Inches(2.0), Inches(0.2),
        "Sources: McKinsey, Gartner, Statista", sz=6, clr=GRAY, font='Inter', align=PP_ALIGN.RIGHT)
    txt(s, CX, Inches(7.0), CW, Inches(1.3),
        "The global AI market has grown from $3.2 billion in 2015 to an estimated $280 billion "
        "in 2026. Companies and individuals who own their AI infrastructure today will capture "
        "a disproportionate share of this value. Bonfire Terminal positions you on the right "
        "side of this curve.",
        sz=9, clr=DARK, font='Inter')

    # ============================================================
    # P7: BENEFITS
    # ============================================================
    print("P7: Benefits")
    s = prs.slides.add_slide(BL); bg(s, WHITE)
    header(s); footer(s, 6)
    hero_bottom = hero_img(s, "cortana-halo-extended", top=HEADER_BOTTOM, max_h=HERO_H)
    txt(s, CX, hero_bottom + Inches(0.1), CW, Inches(0.35), "WHY OWN YOUR AI?",
        sz=18, clr=DARK, bold=True, font='Orbitron')
    bens = [
        ("HARDER", "OpenClaw and cloud chatbots are sandboxed \u2014 they can\u2019t touch "
         "your files, run your software, or control your hardware. OpenClaw requires "
         "you to self-host a Node.js server and burns $400\u2013500/month in hidden API "
         "costs. Bonfire\u2019s compiled daemon runs invisibly on your machine \u2014 bridges "
         "(Telegram, Slack, Discord) and AI providers (Claude, Gemini, Grok) are "
         "built directly in. Zero server setup. It has full access to your GPU, disk, "
         "and installed applications \u2014 for a one-time cost."),
        ("BETTER", "OpenClaw only processes what you route through its plugin system. "
         "Cloud chatbots only know what you paste into a chat window. Bonfire\u2019s "
         "local AI reads your files, understands your project context, and references "
         "your history across sessions. It accesses local documents, spreadsheets, "
         "images, and code \u2014 producing better output because it has better input. "
         "No copy-pasting, no re-explaining, no lost context."),
        ("FASTER", "Every OpenClaw task triggers 3\u20138 API calls at 500ms+ each, "
         "adding 5\u201315 seconds of latency per action. Cloud AI means uploading, "
         "waiting, downloading. Bonfire processes locally at native speed \u2014 50\u2013200ms "
         "first-token response. Whisper transcribes your voice instantly. Just speak, "
         "and your content renders in seconds. No bandwidth limits, no rate limits, "
         "no file-size restrictions."),
        ("STRONGER", "OpenClaw\u2019s skill marketplace had a 12% malware infection "
         "rate \u2014 341 malicious skills including keyloggers and data stealers. 63% of "
         "exposed OpenClaw instances were vulnerable to remote code execution. Cloud AI "
         "platforms score as low as Grade D on privacy audits and train on your data by "
         "default. Bonfire runs 100% locally \u2014 your data physically cannot leave your "
         "machine. Private by architecture, not by promise."),
    ]
    ben_y = hero_bottom + Inches(0.55)
    for t, d in bens:
        txt(s, CX, ben_y, Inches(1.8), Inches(0.3), t,
            sz=14, clr=RED, bold=True, font='Orbitron')
        txt(s, CX+Inches(1.8), ben_y - Inches(0.05), Inches(4.2), Inches(1.1), d,
            sz=8, clr=DARK, font='Inter')
        ben_y += Inches(1.25)

    # ============================================================
    # P8: TYPES OF AI
    # ============================================================
    print("P8: What Bonfire Connects")
    s = prs.slides.add_slide(BL); bg(s, WHITE)
    header(s); footer(s, 7)
    hero_bottom = hero_img(s, "deep-blue", top=HEADER_BOTTOM, max_h=HERO_H)
    txt(s, CX, hero_bottom + Inches(0.1), CW, Inches(0.3), "WHAT BONFIRE CONNECTS",
        sz=15, clr=DARK, bold=True, font='Orbitron')
    txt(s, CX, hero_bottom + Inches(0.45), CW, Inches(0.7),
        "The AI landscape divides into connectivity tools that move data between services, "
        "and content tools that create original output. Most platforms do one or the other. "
        "Bonfire Terminal does both \u2014 locally on your machine with a perpetual license.",
        sz=9, clr=CHAR, font='Inter')
    sep_y = hero_bottom + Inches(1.2)
    rect(s, CX, sep_y, CW, Inches(0.02), GOLD)

    rect(s, CX, sep_y + Inches(0.1), Inches(2.9), Inches(0.3), DKGOLD)
    txt(s, CX, sep_y + Inches(0.12), Inches(2.9), Inches(0.26), "CONNECTIVITY TOOLS",
        sz=10, clr=WHITE, bold=True, font='Orbitron', align=PP_ALIGN.CENTER)
    rect(s, Inches(3.5), sep_y + Inches(0.1), Inches(3.1), Inches(0.3), GOLD)
    txt(s, Inches(3.5), sep_y + Inches(0.12), Inches(3.1), Inches(0.26), "CONTENT TOOLS",
        sz=10, clr=DARK, bold=True, font='Orbitron', align=PP_ALIGN.CENTER)

    txt_y = sep_y + Inches(0.5)
    txt(s, CX, txt_y, Inches(2.9), Inches(1.2),
        "Connectivity platforms like Zapier, Make, and n8n automate "
        "workflows between services. They move data, trigger actions, and "
        "synchronize information. Essential infrastructure, but they "
        "don\u2019t create anything \u2014 they are pipes, not producers.",
        sz=8, clr=CHAR, font='Inter')
    txt(s, Inches(3.5), txt_y, Inches(3.1), Inches(1.2),
        "Content tools like Claude, Gemini, and Whisper "
        "create original output from prompts: text, images, video, audio, "
        "and code. But cloud versions charge monthly, your data lives on their "
        "servers, and they can change pricing or terms at any time.",
        sz=8, clr=CHAR, font='Inter')

    logo_y = txt_y + Inches(1.2)
    conn = ["zapier","slack","discord","mailchimp","wordpress",
            "google-ads","shopify","stripe","meta-color","youtube"]
    content = ["openai","claude-color","gemini-color","xai-grok","meta-llama",
               "elevenlabs","ffmpeg","whisper-openai","stability-ai","blender"]
    logo_row(s, conn[:5], CX+Inches(0.05), logo_y, icon_size=0.4, gap=0.55, max_count=5)
    logo_row(s, conn[5:], CX+Inches(0.05), logo_y + Inches(0.5), icon_size=0.4, gap=0.55, max_count=5)
    logo_row(s, content[:5], Inches(3.55), logo_y, icon_size=0.4, gap=0.55, max_count=5)
    logo_row(s, content[5:], Inches(3.55), logo_y + Inches(0.5), icon_size=0.4, gap=0.55, max_count=5)

    box_y = logo_y + Inches(1.1)
    rect(s, CX, box_y, CW, Inches(1.2), DARK)
    txt(s, CX+Inches(0.15), box_y + Inches(0.1), Inches(5.9), Inches(1.0),
        "Bonfire Terminal sits at the intersection of both categories:\n\n"
        "It CONNECTS to any AI provider \u2014 OpenAI, Claude, Gemini, Grok, or local models.\n"
        "It CREATES content across every medium \u2014 video, transcription, images, text, code.\n"
        "It runs LOCALLY with zero recurring fees and full data privacy.",
        sz=8, clr=GOLD, font='Inter', align=PP_ALIGN.CENTER)

    # ============================================================
    # P9: PRODUCT CATALOG GRID — Section 2
    # ============================================================
    print("P9: Product Catalog Grid")
    s = prs.slides.add_slide(BL); bg(s, WHITE)
    header(s); footer(s, 8)
    hero_bottom = hero_img(s, "jarvis-hologram", top=HEADER_BOTTOM, max_h=HERO_H)
    txt(s, CX, hero_bottom + Inches(0.05), Inches(5.8), Inches(0.3),
        "THE BONFIRE PLATFORM", sz=16, clr=DARK, bold=True, font='Orbitron')

    grid_y = hero_bottom + Inches(0.4)
    ai_providers = ["claude-color","openai","gemini-color","xai-grok","meta-llama"]
    logo_grid_row(s, "AI PROVIDERS",
                  "Claude, Codex, Gemini, Grok + Llama (offline). Switch models instantly.",
                  ai_providers, CX, grid_y, CW, Inches(1.4))

    aff_nets = ["clickbank","digistore24","impact","cj-affiliate",
                "awin","rakuten","jvzoo","walmart"]
    logo_grid_row(s, "AFFILIATE NETWORKS",
                  "Connect to 16 affiliate platforms. Track commissions. Automate campaigns.",
                  aff_nets, CX, grid_y + Inches(1.6), CW, Inches(1.4))

    # Messaging bridges row
    bridges = ["telegram","discord","slack","whatsapp","twilio","zapier"]
    logo_grid_row(s, "MESSAGING BRIDGES",
                  "Telegram, Discord, Slack, WhatsApp, Twilio, Zapier \u2014 compiled into the daemon.",
                  bridges, CX, grid_y + Inches(3.2), CW, Inches(1.4))

    # ============================================================
    # P10: MORE PRODUCTS — V6: Added page title
    # ============================================================
    print("P10: Learning Roadmap")
    s = prs.slides.add_slide(BL); bg(s, WHITE)
    header(s); footer(s, 9)
    hero_bottom = hero_img(s, "cyberpunk-johnny", top=HEADER_BOTTOM, max_h=HERO_H)

    txt(s, CX, hero_bottom + Inches(0.05), CW, Inches(0.3),
        "YOUR BONFIRE LEARNING PATH", sz=14, clr=DARK, bold=True, font='Orbitron')
    txt(s, CX, hero_bottom + Inches(0.35), CW, Inches(0.2),
        "A proven 5-level system from setup to scaling", sz=8, clr=CHAR, font='Inter')

    levels = [
        ("1", "SETUP", "AI tools, affiliate networks, website, ad platforms",
         "Configure Bonfire, connect providers, launch AlphaChat"),
        ("2", "SIGNS OF LIFE", "Optimize targeting, creative, landing pages, offers",
         "Use AI to test and refine \u2014 find what converts"),
        ("3", "GROWING", "Launch campaigns, expand audiences, add creatives",
         "Scale what works across channels with AI automation"),
        ("4", "SCALING", "Split-test, increase budgets, optimize payouts",
         "Multiply revenue with data-driven AI decisions"),
        ("5", "MAINTENANCE", "Automate, build a team, compound profits",
         "Your AI ghost agency runs while you focus on strategy"),
    ]
    lv_y = hero_bottom + Inches(0.65)
    for num, title, desc, bonfire_note in levels:
        # Level number circle
        oval(s, CX, lv_y, Inches(0.45), Inches(0.45), DARK)
        txt(s, CX+Inches(0.02), lv_y+Inches(0.03), Inches(0.41), Inches(0.39),
            num, sz=16, clr=GOLD, bold=True, font='Orbitron', align=PP_ALIGN.CENTER)
        # Connector line
        if num != "5":
            rect(s, CX+Inches(0.2), lv_y+Inches(0.45), Inches(0.04), Inches(0.55), GOLD)
        # Title and description
        txt(s, CX+Inches(0.55), lv_y+Inches(0.0), Inches(2.5), Inches(0.22),
            title, sz=10, clr=DARK, bold=True, font='Orbitron')
        txt(s, CX+Inches(0.55), lv_y+Inches(0.22), Inches(2.5), Inches(0.25),
            desc, sz=7, clr=CHAR, font='Inter')
        # Bonfire note on right
        txt(s, Inches(3.8), lv_y+Inches(0.05), Inches(2.8), Inches(0.4),
            bonfire_note, sz=7, clr=DKGOLD, font='Inter', italic=True)
        lv_y += Inches(0.95)

    # ============================================================
    # P11: CLOUD VS OWNED
    # ============================================================
    print("P11: Cloud vs Owned")
    s = prs.slides.add_slide(BL); bg(s, WHITE)
    header(s); footer(s, 10)
    hero_bottom = hero_img(s, "sonny-irobot", top=HEADER_BOTTOM, max_h=HERO_H)
    rect(s, CX, hero_bottom, CW, Inches(0.8), DARK)
    txt(s, CX+Inches(0.1), hero_bottom + Inches(0.05), Inches(5.9), Inches(0.3),
        "CLOUD AI vs. OWNED AI", sz=16, clr=GOLD, bold=True, font='Orbitron')
    txt(s, CX+Inches(0.1), hero_bottom + Inches(0.4), Inches(5.9), Inches(0.35),
        "The cloud AI model extracts recurring revenue. You pay monthly, own nothing. "
        "Here is how Bonfire Terminal compares:", sz=8, clr=LGOLD, font='Inter')

    tbl_top = hero_bottom + Inches(0.9)
    cxs = [0.4, 2.2, 4.4]; cws = [1.7, 2.1, 2.2]
    rect(s, Inches(cxs[0]), tbl_top, Inches(cws[0]), Inches(0.35), CHAR)
    rect(s, Inches(cxs[1]), tbl_top, Inches(cws[1]), Inches(0.35), GRAY)
    rect(s, Inches(cxs[2]), tbl_top, Inches(cws[2]), Inches(0.35), GOLD)
    txt(s, Inches(cxs[0]), tbl_top+Inches(0.02), Inches(cws[0]), Inches(0.3), "FEATURE",
        sz=8, clr=WHITE, bold=True, font='Inter', align=PP_ALIGN.CENTER)
    txt(s, Inches(cxs[1]), tbl_top+Inches(0.02), Inches(cws[1]), Inches(0.3), "CLOUD AI",
        sz=8, clr=WHITE, bold=True, font='Inter', align=PP_ALIGN.CENTER)
    txt(s, Inches(cxs[2]), tbl_top+Inches(0.02), Inches(cws[2]), Inches(0.3), "BONFIRE TERMINAL",
        sz=8, clr=DARK, bold=True, font='Inter', align=PP_ALIGN.CENTER)
    rows = [
        ["Cost","$20-200/mo per tool\nper user, forever","One-time $5,000\nOwn forever"],
        ["Data Privacy","Your data on their servers\nUsed for model training","100% local processing\nData never leaves"],
        ["Availability","Internet required\nOutages affect work","Works fully offline\nNo dependencies"],
        ["Customization","Limited to their UI\nNo local integration","Full CLI + AlphaChat\nPlugin ecosystem"],
        ["Ownership","You rent monthly\nThey own everything","You OWN the software\nPerpetual license"],
    ]
    y = tbl_top + Inches(0.4)
    row_h = Inches(0.42)
    for i, row in enumerate(rows):
        bgc = OFFWH if i%2==0 else WHITE; gbg = GOLDBG if i%2==0 else WHITE
        rect(s, Inches(cxs[0]), y, Inches(cws[0]), row_h, bgc)
        rect(s, Inches(cxs[1]), y, Inches(cws[1]), row_h, bgc)
        rect(s, Inches(cxs[2]), y, Inches(cws[2]), row_h, gbg)
        txt(s, Inches(cxs[0]+0.05), y+Inches(0.03), Inches(cws[0]-0.1), Inches(0.44),
            row[0], sz=6, clr=DARK, bold=True, font='Inter')
        txt(s, Inches(cxs[1]+0.05), y+Inches(0.03), Inches(cws[1]-0.1), Inches(0.44),
            row[1], sz=5, clr=CHAR, font='Inter')
        txt(s, Inches(cxs[2]+0.05), y+Inches(0.03), Inches(cws[2]-0.1), Inches(0.44),
            row[2], sz=5, clr=DARK, bold=True, font='Inter')
        y += row_h

    # Savings table moved here from P12
    sav_y = y + Inches(0.15)
    txt(s, CX, sav_y, CW, Inches(0.2),
        "CLOUD AI SUBSCRIPTION COST OVER TIME", sz=8, clr=GOLD, bold=True, font='Inter')
    tbl_y = sav_y + Inches(0.25)
    cx_cols = [0.4,1.5,2.5,3.5,4.5,5.5]
    rect(s, CX, tbl_y, CW, Inches(0.22), DARK)
    cols = ["","Year 1","Year 2","Year 3","Year 5","Year 10"]
    for i,c in enumerate(cols):
        txt(s, Inches(cx_cols[i]), tbl_y+Inches(0.01), Inches(0.95), Inches(0.2), c,
            sz=6, clr=WHITE, bold=True, font='Inter', align=PP_ALIGN.CENTER)
    rect(s, CX, tbl_y+Inches(0.25), CW, Inches(0.28), OFFWH)
    for i,v in enumerate(["Cloud AI","$2,400","$4,800","$7,200","$12,000","$24,000"]):
        txt(s, Inches(cx_cols[i]), tbl_y+Inches(0.26), Inches(0.95), Inches(0.26), v,
            sz=6, clr=CHAR, font='Inter', align=PP_ALIGN.CENTER)
    rect(s, CX, tbl_y+Inches(0.56), CW, Inches(0.28), GOLDBG)
    for i,v in enumerate(["Bonfire","$5,000","$5,000","$5,000","$5,000","$5,000"]):
        txt(s, Inches(cx_cols[i]), tbl_y+Inches(0.57), Inches(0.95), Inches(0.26), v,
            sz=6, clr=DKGOLD if i>0 else DARK, bold=True, font='Inter', align=PP_ALIGN.CENTER)
    rect(s, CX, tbl_y+Inches(0.87), CW, Inches(0.28), DARK)
    for i,v in enumerate(["SAVINGS","-$2,600","-$200","+$2,200","+$7,000","+$19,000"]):
        c = RGBColor(0xFF,0x66,0x55) if v.startswith("-") else GREEN
        if i==0: c=WHITE
        txt(s, Inches(cx_cols[i]), tbl_y+Inches(0.88), Inches(0.95), Inches(0.26), v,
            sz=6, clr=c, bold=True, font='Inter', align=PP_ALIGN.CENTER)

    # ============================================================
    # P12: ROI, TESTIMONIALS & AIDE INCOME
    # ============================================================
    print("P12: ROI & Testimonials")
    s = prs.slides.add_slide(BL); bg(s, WHITE)
    header(s); footer(s, 11)
    hero_bottom = hero_img(s, "pdp-11", top=HEADER_BOTTOM, max_h=HERO_H)
    rect(s, CX, hero_bottom, CW, Inches(0.45), DARK)
    txt(s, CX+Inches(0.1), hero_bottom + Inches(0.05), Inches(5.9), Inches(0.35),
        "RESULTS & INCOME POTENTIAL", sz=14, clr=GOLD, bold=True, font='Orbitron')

    # Student testimonials
    test_y = hero_bottom + Inches(0.55)
    txt(s, CX, test_y, CW, Inches(0.2),
        "WHAT OUR STUDENTS SAY", sz=9, clr=GOLD, bold=True, font='Inter')
    testimonials = [
        ("\u201CWe made 150 sales and our first $1,114.20 in our first month.\u201D",
         "\u2014 Current Student, AI Marketing Mentorship"),
        ("\u201CI got my 6th commission. All these results are 100% from implementing "
         "what was taught. Super excited... now just working to get bigger numbers.\u201D",
         "\u2014 Current Student, AIMS Graduate"),
        ("\u201CI wanted good information by people who knew what they were talking "
         "about and are good at what they do.\u201D",
         "\u2014 Amy, Telehealth Niche"),
    ]
    ty = test_y + Inches(0.3)
    for quote, attr in testimonials:
        rect(s, CX, ty, CW, Inches(0.75), OFFWH, line=GOLD)
        txt(s, CX+Inches(0.15), ty+Inches(0.05), Inches(5.8), Inches(0.45),
            quote, sz=7, clr=DARK, font='Inter', italic=True)
        txt(s, CX+Inches(0.15), ty+Inches(0.5), Inches(5.8), Inches(0.2),
            attr, sz=6, clr=DKGOLD, bold=True, font='Inter')
        ty += Inches(0.85)

    # AIDE Income section
    rect(s, CX, ty + Inches(0.1), CW, Inches(0.25), GOLD)
    txt(s, CX+Inches(0.1), ty + Inches(0.12), Inches(5.9), Inches(0.21),
        "AIDE INCOME POTENTIAL", sz=9, clr=DARK, bold=True, font='Orbitron')
    comm_y = ty + Inches(0.4)
    rect(s, CX, comm_y, Inches(3.0), Inches(0.5), DARK)
    txt(s, CX+Inches(0.1), comm_y+Inches(0.05), Inches(2.8), Inches(0.2),
        "90%", sz=18, clr=GOLD, bold=True, font='Orbitron')
    txt(s, CX+Inches(0.1), comm_y+Inches(0.3), Inches(2.8), Inches(0.15),
        "You Keep  |  Only 10% Licensing Fee", sz=6, clr=LGOLD, font='Inter')
    rect(s, Inches(3.6), comm_y, Inches(3.0), Inches(0.5), DARK)
    txt(s, Inches(3.7), comm_y+Inches(0.05), Inches(2.8), Inches(0.2),
        "80%", sz=18, clr=GOLD, bold=True, font='Orbitron')
    txt(s, Inches(3.7), comm_y+Inches(0.3), Inches(2.8), Inches(0.15),
        "You Keep  |  Only 20% on Renewals", sz=6, clr=LGOLD, font='Inter')

    inc_y = comm_y + Inches(0.6)
    inc_data = [
        ("1 client/mo", "$54,000/yr"),
        ("2.5 clients/mo", "$135,000/yr"),
        ("5 clients/mo", "$270,000/yr"),
    ]
    col_w = Inches(2.06)
    for i, (label, amount) in enumerate(inc_data):
        x = CX + Inches(i * 2.07)
        bgc = GOLDBG if i < 2 else DARK
        txt_c = DARK if i < 2 else GOLD
        lbl_c = DKGOLD if i < 2 else LGOLD
        rect(s, x, inc_y, col_w, Inches(0.6), bgc, line=GOLD)
        txt(s, x+Inches(0.05), inc_y+Inches(0.05), Inches(1.96), Inches(0.3),
            amount, sz=14, clr=txt_c, bold=True, font='Orbitron', align=PP_ALIGN.CENTER)
        txt(s, x+Inches(0.05), inc_y+Inches(0.38), Inches(1.96), Inches(0.18),
            label, sz=7, clr=lbl_c, font='Inter', align=PP_ALIGN.CENTER)

    # ============================================================
    # P13: WHO BENEFITS
    # ============================================================
    print("P13: Who Benefits")
    s = prs.slides.add_slide(BL); bg(s, WHITE)
    header(s); footer(s, 12)
    hero_bottom = hero_img(s, "jobs-macintosh", top=HEADER_BOTTOM, max_h=HERO_H)
    rect(s, CX, hero_bottom, CW, Inches(0.45), DARK)
    txt(s, CX+Inches(0.1), hero_bottom + Inches(0.05), Inches(5.9), Inches(0.35),
        "WHO IS BONFIRE FOR?", sz=16, clr=GOLD, bold=True, font='Orbitron')
    personas = [
        ("CONTENT CREATORS", "YouTubers, podcasters, and social media managers producing "
         "content at scale. AI-powered video editing, transcription, and scriptwriting."),
        ("AFFILIATE MARKETERS", "Run ads, build funnels, manage campaigns across 16 affiliate "
         "platforms. Automate content creation for every channel."),
        ("SMALL BUSINESS OWNERS", "Professional marketing output without a marketing department. "
         "Generate ad copy, email sequences, social posts, and video."),
        ("IT CONSULTANTS & AGENCIES", "Deploy Bonfire for clients as a Certified AIDE. "
         "Keep 90% of every sale \u2014 only 10% licensing fee to us."),
        ("FREELANCERS", "10x your content output while maintaining full creative control. "
         "Serve more clients in less time with AI tools at native speed."),
        ("ENTREPRENEURS", "AI infrastructure you own outright. Scale without escalating "
         "cloud costs. Build on a foundation you control."),
    ]
    y = hero_bottom + Inches(0.55)
    for i, (t, d) in enumerate(personas):
        color = GOLD if i%2==0 else NAVY
        rect(s, CX, y, CW, Inches(0.72), WHITE)
        rect(s, CX, y + Inches(0.02), Inches(0.06), Inches(0.68), color)
        txt(s, CX+Inches(0.15), y + Inches(0.02), Inches(5.8), Inches(0.2), t,
            sz=8, clr=DARK, bold=True, font='Inter')
        txt(s, CX+Inches(0.15), y + Inches(0.25), Inches(5.8), Inches(0.4), d,
            sz=7, clr=CHAR, font='Inter')
        y += Inches(0.78)
    rect(s, CX, Inches(8.0), CW, Inches(0.5), DARK)
    txt(s, CX+Inches(0.2), Inches(8.05), Inches(5.8), Inches(0.4),
        "\u201CThe best time to own your AI was yesterday. The second best time is today.\u201D",
        sz=10, clr=GOLD, font='Inter', italic=True, align=PP_ALIGN.CENTER)

    # ============================================================
    # P14: TECH SPECS
    # ============================================================
    print("P14: Technical Specs")
    s = prs.slides.add_slide(BL); bg(s, WHITE)
    header(s); footer(s, 13)
    hero_bottom = hero_img(s, "univac-computer", top=HEADER_BOTTOM, max_h=HERO_H)
    rect(s, CX, hero_bottom, CW, Inches(0.7), DARK)
    txt(s, CX+Inches(0.1), hero_bottom + Inches(0.05), Inches(5.9), Inches(0.3),
        "WHAT\u2019S UNDER THE HOOD", sz=15, clr=GOLD, bold=True, font='Orbitron')
    txt(s, CX+Inches(0.1), hero_bottom + Inches(0.4), Inches(5.9), Inches(0.25),
        "50+ tools and platforms across four categories. All local, no cloud middleware.",
        sz=8, clr=LGOLD, font='Inter')

    cats = [
        ("AI PROVIDERS", DKGOLD,
         ["claude-color","gemini-color","xai-grok","openai","meta-llama"],
         "Claude, Gemini, Grok, OpenAI + offline Llama models."),
        ("VIDEO & AUDIO", DARK,
         ["ffmpeg","whisper-openai","elevenlabs","blender","remotion","kdenlive"],
         "FFmpeg, Whisper, ElevenLabs, Blender, Remotion, Kdenlive."),
        ("DESIGN & CONTENT", DKGOLD,
         ["stability-ai","canva","davinci-resolve","adobe-premiere-pro",
          "figma","imagemagick"],
         "Image generation, design, video editing, post-production."),
        ("MARKETING", DARK,
         ["youtube","instagram","tiktok","shopify","stripe",
          "mailchimp","google-ads","wordpress","meta-color","discord"],
         "Publish to every platform. Manage ads and e-commerce."),
    ]
    y = hero_bottom + Inches(0.8)
    for cat, color, cat_logos, desc in cats:
        rect(s, CX, y, CW, Inches(0.22), color)
        txt(s, CX+Inches(0.08), y+Inches(0.01), Inches(3.0), Inches(0.2), cat,
            sz=7, clr=WHITE, bold=True, font='Inter')
        txt(s, Inches(3.5), y+Inches(0.01), Inches(3.0), Inches(0.2), desc,
            sz=5, clr=LGOLD, font='Inter')
        logo_row(s, cat_logos, CX+Inches(0.02), y + Inches(0.26),
                 icon_size=0.3, gap=0.6, max_count=10)
        y += Inches(0.85)
    txt(s, CX, Inches(8.1), CW, Inches(0.25),
        "Requirements: Windows 10+ / macOS 12+  |  8 GB RAM  |  10 GB disk",
        sz=7, clr=CHAR, font='Inter', align=PP_ALIGN.CENTER)

    # ============================================================
    # P15: PRICING — V6: Redesigned layout
    # $5k full-width top, $10k/$18k side-by-side, JC McLaren photo
    # ============================================================
    print("P15: Pricing")
    s = prs.slides.add_slide(BL); bg(s, WHITE)
    header(s); footer(s, 14)
    txt(s, CX, Inches(0.7), CW, Inches(0.4),
        "CHOOSE YOUR EDITION", sz=16, clr=DARK, bold=True, font='Orbitron')

    # Helper for green checkmark bullet lines
    def check_bullets(s, x, y, items, sz=7):
        for item in items:
            txt(s, x, y, Inches(0.2), Inches(0.16),
                "\u2713", sz=sz, clr=CHECKGREEN, bold=True, font='Inter')
            txt(s, x+Inches(0.2), y, Inches(5.0), Inches(0.16),
                item, sz=sz, clr=CHAR, font='Inter')
            y += Inches(0.17)
        return y

    # Card 1: Bonfire Terminal — FULL WIDTH
    rect(s, CX, Inches(1.15), CW, Inches(2.8), OFFWH, line=DARK)
    rect(s, CX, Inches(1.15), CW, Inches(0.45), DARK)
    txt(s, CX+Inches(0.15), Inches(1.17), Inches(5.9), Inches(0.22),
        "BONFIRE TERMINAL", sz=13, clr=GOLD, bold=True, font='Orbitron')
    txt(s, CX+Inches(0.15), Inches(1.40), Inches(5.9), Inches(0.18),
        "The Complete AI Ownership Platform", sz=7, clr=LGOLD, font='Inter')
    txt(s, CX+Inches(0.15), Inches(1.65), Inches(3.0), Inches(0.4),
        "$5,000", sz=36, clr=RED, bold=True, font='Orbitron')
    txt(s, CX+Inches(2.8), Inches(1.72), Inches(1.0), Inches(0.25),
        "ONE TIME", sz=8, clr=GOLD, bold=True, font='Inter')
    check_bullets(s, CX+Inches(0.15), Inches(2.1), [
        "Perpetual license \u2014 buy once, own forever",
        "5 AI providers: Claude, Codex, Gemini, Grok + offline Llama",
        "FFmpeg video, Whisper speech-to-text, AlphaChat terminal",
        "Plugin ecosystem + offline capability with TinyLlama",
        "AI Marketers Lifetime Club + 12-module AIMS course",
        "1 onboarding session + 1 AI coaching session",
        "20% affiliate commission on all referrals",
    ], sz=7)

    # Card 2: + Coaching — LEFT HALF (equalized widths)
    half_w = Inches(3.0)
    card2_top = Inches(4.15)
    rect(s, CX, card2_top, half_w, Inches(2.0), WHITE, line=DARK)
    rect(s, CX, card2_top, half_w, Inches(0.35), DARK)
    txt(s, CX+Inches(0.1), card2_top+Inches(0.04), Inches(2.8), Inches(0.27),
        "+ COACHING", sz=11, clr=WHITE, bold=True, font='Orbitron')
    txt(s, CX+Inches(0.1), card2_top+Inches(0.42), Inches(2.8), Inches(0.35),
        "$10,000", sz=24, clr=DARK, bold=True, font='Orbitron')
    check_bullets(s, CX+Inches(0.1), card2_top+Inches(0.8), [
        "Everything in Bonfire Terminal",
        "8 private 1-on-1 coaching sessions",
        "Workflow setup + prompt engineering",
        "Campaign strategy + monthly calls",
        "Direct Slack access to the team",
    ], sz=6)

    # Card 3: + AIDE Certification — RIGHT HALF (equalized widths)
    card3_x = CX + half_w + Inches(0.2)
    card3_w = CW - half_w - Inches(0.2)
    rect(s, card3_x, card2_top, card3_w, Inches(2.0), WHITE, line=GOLD)
    rect(s, card3_x, card2_top, card3_w, Inches(0.35), GOLD)
    txt(s, card3_x+Inches(0.1), card2_top+Inches(0.04), Inches(2.7), Inches(0.27),
        "+ AIDE CERTIFICATION", sz=10, clr=DARK, bold=True, font='Orbitron')
    txt(s, card3_x+Inches(0.1), card2_top+Inches(0.42), Inches(2.7), Inches(0.35),
        "$18,000", sz=24, clr=DARK, bold=True, font='Orbitron')
    check_bullets(s, card3_x+Inches(0.1), card2_top+Inches(0.8), [
        "Everything in Coaching (8 sessions)",
        "10 AIDE certification trainings",
        "Reseller rights (keep 90% of sales)",
        "80% recurring on support renewals",
        "Sales materials + dedicated support",
    ], sz=6)

    # JC McLaren photo — with gap above cards
    jc_section_top = Inches(6.3)
    jc_h = Inches(2.0)
    rect(s, CX, jc_section_top, CW, jc_h, DARK)
    if os.path.exists(JC_MCLAREN):
        pic(s, JC_MCLAREN, CX, jc_section_top, w=Inches(3.6), h=jc_h)
    txt(s, CX+Inches(3.7), jc_section_top+Inches(0.2), Inches(2.4), Inches(0.3),
        "Schedule a call today", sz=12, clr=GOLD, bold=True, font='Orbitron')
    txt(s, CX+Inches(3.7), jc_section_top+Inches(0.6), Inches(2.4), Inches(0.3),
        "bonfireterminal.com/apply", sz=10, clr=WHITE, font='Inter')
    txt(s, CX+Inches(3.7), jc_section_top+Inches(1.1), Inches(2.4), Inches(0.8),
        "All packages include a perpetual license.\n\n"
        "No subscriptions.\nNo recurring fees.",
        sz=9, clr=WHITE, font='Inter')

    # ============================================================
    # P16: SOCIAL PROOF — V6: Fixed logo overlap, updated section title
    # ============================================================
    print("P16: Social Proof")
    s = prs.slides.add_slide(BL); bg(s, WHITE)
    header(s); footer(s, 15)
    txt(s, CX, Inches(0.7), Inches(5.8), Inches(0.35),
        "WHY CHOOSE BONFIRE TERMINAL?", sz=15, clr=DARK, bold=True, font='Orbitron')

    # Trust badges — fixed alignment
    aide_badge = generate_aide_badge()
    badge_files = []
    for sub, pat in [("certification","g2-high-performer"),
                     ("certification","ssl-secured"),
                     ("certification","apple-developer"),
                     ("certification","clickbank-platinum")]:
        b = img(sub, pat)
        if b: badge_files.append(b)
    # AIDE badge removed per user request
    badge_h = Inches(0.55)
    total_badges = len(badge_files)
    if total_badges > 0:
        # Last badge (AIDE) gets extra right offset
        spacing = 5.5 / max(total_badges - 1, 1)
        for i, bf in enumerate(badge_files):
            x = Inches(0.5 + i * spacing)
            if i == total_badges - 1:
                x = Inches(5.8)  # Push AIDE badge to far right
            pic(s, bf, x, Inches(1.2), h=badge_h)

    txt(s, CX, Inches(1.9), CW, Inches(0.22),
        "WHAT INDUSTRY LEADERS SAY ABOUT AI", sz=9, clr=GOLD, bold=True, font='Inter')

    # 3 testimonials
    tqs = [
        ("elon-musk", "Elon Musk",
         "\u201CThe most important development in the history of technology. "
         "AI will be more profound than electricity or fire.\u201D",
         "CEO of Tesla, SpaceX, and xAI"),
        ("jensen-huang", "Jensen Huang",
         "\u201CAI is the most powerful technology force of our time. Every industry "
         "will be revolutionized by artificial intelligence.\u201D",
         "CEO of NVIDIA"),
        ("mark-cuban", "Mark Cuban",
         "\u201CThe world\u2019s first trillionaires are going to come from somebody who "
         "masters AI and all its derivatives.\u201D",
         "Investor, Entrepreneur, Shark Tank"),
    ]
    y = 2.3
    for ik, nm, q, title in tqs:
        circ_img = str(CIRC / f"{ik}-circle.png")
        if os.path.exists(circ_img):
            oval(s, CX, Inches(y+0.05), Inches(0.8), Inches(0.8), WHITE, line=GOLD)
            pic(s, circ_img, CX+Inches(0.05), Inches(y+0.1), w=Inches(0.7), h=Inches(0.7))
        txt(s, CX+Inches(0.95), Inches(y), Inches(5.0), Inches(0.5), q,
            sz=9, clr=DARK, font='Inter', italic=True)
        txt(s, CX+Inches(0.95), Inches(y+0.55), Inches(5.0), Inches(0.2),
            f"\u2014 {nm}, {title}", sz=8, clr=GOLD, bold=True, font='Inter')
        if y < 4.2:
            rect(s, CX+Inches(0.95), Inches(y+0.85), Inches(5.0), Inches(0.01), LGOLD)
        y += 1.0

    # Customer-facing metrics
    rect(s, CX, Inches(5.5), CW, Inches(0.6), DARK)
    txt(s, CX+Inches(0.1), Inches(5.55), Inches(6.0), Inches(0.5),
        "5 AI Providers  |  50+ Integrations  |  Perpetual License\n"
        "Zero Recurring Fees  |  100% Offline Capable  |  Own Your Data",
        sz=9, clr=GOLD, font='Inter', align=PP_ALIGN.CENTER, bold=True)

    # John Crestani + media
    rect(s, CX, Inches(6.3), CW, Inches(0.35), GOLD)
    txt(s, CX, Inches(6.32), CW, Inches(0.3),
        "JOHN CRESTANI: FEATURED ON 100+ TV STATIONS WORLDWIDE",
        sz=8, clr=DARK, bold=True, font='Inter', align=PP_ALIGN.CENTER)
    txt(s, CX, Inches(6.75), CW, Inches(0.2), "AS SEEN ON",
        sz=8, clr=CHAR, bold=True, font='Inter', align=PP_ALIGN.CENTER)

    # V6 FIX: Only 5 logos, constrained width to prevent overlap
    seen = ['forbes','cnbc','bloomberg','fox-business','wired']
    logo_w = Inches(0.9)  # Fixed width for all logos
    total_logos = len(seen)
    logo_area = 6.0  # Available width in inches
    logo_gap = logo_area / total_logos
    for i, ln in enumerate(seen):
        li = img("as-seen-on", ln)
        x = Inches(0.5 + i * logo_gap)
        if li:
            pic(s, li, x, Inches(7.0), w=logo_w)
        else:
            txt(s, x, Inches(7.0), logo_w, Inches(0.4),
                ln.upper().replace("-"," "), sz=6, clr=CHAR, font='Inter', align=PP_ALIGN.CENTER)

    txt(s, CX, Inches(7.55), CW, Inches(0.5),
        "John Crestani has been featured in Forbes, Entrepreneur, CNBC, Bloomberg, "
        "and over 100 television stations worldwide.",
        sz=7, clr=CHAR, font='Inter', align=PP_ALIGN.CENTER)

    # ============================================================
    # P17: FAQ — V6: Updated OpenClaw FAQ, redesigned JC section
    # ============================================================
    print("P17: FAQ")
    s = prs.slides.add_slide(BL); bg(s, WHITE)
    header(s); footer(s, 16)
    txt(s, CX, Inches(0.7), CW, Inches(0.4), "FREQUENTLY ASKED QUESTIONS",
        sz=16, clr=DARK, bold=True, font='Orbitron')
    faqs = [
        ("Do I actually own the software?",
         "Yes. Perpetual license. No subscriptions, no recurring fees. Buy once, own forever."),
        ("Does it work without internet?",
         "Yes. All media processing and local AI models work fully offline."),
        ("How is Bonfire different from OpenClaw?",
         "Bonfire processes locally, works offline, and is built for marketers \u2014 not developers."),
        ("Can I use this for my clients?",
         "Yes. AIDE certification gives you reseller rights \u2014 keep 90%, pay only 10% licensing fee."),
        ("Can I reschedule coaching calls?",
         "Yes, on a discretionary basis. We respect your time and ask that you "
         "respect your coach\u2019s time in return. Rescheduling is available with "
         "reasonable notice."),
        ("What support do I get?",
         "Community access, weekly webinars, email support. Coaching adds priority + Slack."),
    ]
    col_w = Inches(2.9)
    for i, (q, a) in enumerate(faqs):
        col = i % 2; row = i // 2
        x = CX + Inches(col * 3.2)
        y = Inches(1.3 + row * 1.15)
        txt(s, x, y, col_w, Inches(0.22), q, sz=9, clr=RED, bold=True, font='Inter')
        txt(s, x, y+Inches(0.25), col_w, Inches(0.7), a, sz=8, clr=DARK, font='Inter')
        rect(s, x, y+Inches(0.95), Inches(2.6), Inches(0.005), LGOLD)

    # V6.2: JC personal statement — expanded, larger image
    jc_y = Inches(5.2)
    jc_h = Inches(3.1)
    rect(s, CX, jc_y, CW, jc_h, DARK)
    if os.path.exists(JC_STUDENT):
        pic(s, JC_STUDENT, CX, jc_y, w=Inches(3.8), h=jc_h)
    txt(s, CX+Inches(3.9), jc_y+Inches(0.15), Inches(2.2), Inches(0.3),
        "A NOTE FROM JOHN CRESTANI", sz=9, clr=GOLD, bold=True, font='Orbitron')
    txt(s, CX+Inches(3.9), jc_y+Inches(0.55), Inches(2.2), Inches(1.6),
        "I built Bonfire Terminal because I believe everyone deserves to own their AI tools, "
        "not rent them.\n\n"
        "My commitment is simple: help you build real, lasting wealth online "
        "using the most powerful technology ever created.\n\n"
        "Every feature, every integration, "
        "every update is designed with one goal \u2014 your success.",
        sz=8, clr=WHITE, font='Inter', italic=True)
    if os.path.exists(JC_SIG_WHITE):
        pic(s, JC_SIG_WHITE, CX+Inches(3.9), jc_y+Inches(2.3), w=Inches(1.8))

    # ============================================================
    # P18: BACK COVER
    # ============================================================
    print("P18: Back Cover")
    s = prs.slides.add_slide(BL); bg(s, DARK)
    rect(s, Inches(0.25), Inches(0.25), Inches(6.5), Inches(8.5), DARK, line=GOLD)
    rect(s, Inches(0.35), Inches(0.35), Inches(6.3), Inches(8.3),
         DARK, line=RGBColor(0x3A,0x35,0x28))
    if os.path.exists(LOGO_LG):
        pic(s, LOGO_LG, Inches(2.75), Inches(1.5), w=Inches(1.5), h=Inches(1.5))
    txt(s, Inches(0.5), Inches(3.2), Inches(6), Inches(0.6), "BONFIRE TERMINAL",
        sz=32, clr=WHITE, bold=True, font='Orbitron', align=PP_ALIGN.CENTER)
    rect(s, Inches(2.5), Inches(4.0), Inches(2), Inches(0.03), GOLD)
    qr_path = str(IMG / "qr-code-apply.png")
    if os.path.exists(qr_path):
        rect(s, Inches(2.6), Inches(4.9), Inches(1.8), Inches(1.8), WHITE, line=GOLD)
        pic(s, qr_path, Inches(2.7), Inches(5.0), w=Inches(1.6), h=Inches(1.6))
    txt(s, Inches(0.5), Inches(6.9), Inches(6), Inches(0.3),
        "Scan to apply  |  bonfireterminal.com/apply",
        sz=10, clr=GOLD, font='Inter', align=PP_ALIGN.CENTER)
    txt(s, Inches(0.5), Inches(7.5), Inches(6), Inches(0.3),
        "bonfireterminal.com",
        sz=10, clr=GOLD, font='Inter', align=PP_ALIGN.CENTER)

    print(f"\nSaving {OUT}...")
    prs.save(str(OUT))
    print(f"DONE! {len(prs.slides)} slides saved to {OUT}")


if __name__ == "__main__":
    build()
